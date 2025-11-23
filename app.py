import originpro as op
import pandas as pd
import numpy as np
import os
import sys
import threading
import webview
import pythoncom
import time
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from datetime import date, datetime
from pptx import Presentation
from pptx.util import Inches

# --- Path Setup ---
if getattr(sys, 'frozen', False):
    base_dir = sys._MEIPASS
    template_folder = os.path.join(base_dir, 'react_build')
    static_folder = os.path.join(base_dir, 'react_build', 'static')
else:
    base_dir = os.getcwd()
    template_folder = 'react_build'
    static_folder = 'react_build/static'

app = Flask(__name__, static_folder=template_folder, static_url_path='')
CORS(app) # Fixes the "Failed" false alarm

# --- Logging & Context ---
def log_status(message):
    print(f"[{datetime.now().strftime('%H:%M:%S')}] {message}")

class OriginContext:
    def __enter__(self):
        pythoncom.CoInitialize()
        op.set_show(True)
        return self
    def __exit__(self, exc_type, exc_value, traceback):
        pythoncom.CoUninitialize()

# --- Data Loaders ---
def get_delimiter(file_obj):
    try:
        sample = file_obj.read(1024).decode('latin1', errors='ignore') 
        file_obj.seek(0)
        if ',' in sample: return ','
        if '\t' in sample: return '\t'
        return ','
    except Exception: return ','

def find_data_start_row(file_obj):
    file_obj.seek(0)
    try:
        lines = file_obj.read().decode('latin1').splitlines()
    except Exception: return 0
    file_obj.seek(0)
    for i, line in enumerate(lines):
        if '[Data]' in line: return i + 1
    return 0

def load_data(file_obj, skiprows, usecols, colnames):
    delim = get_delimiter(file_obj)
    file_obj.seek(0) 
    if skiprows == -1:
        skiprows = find_data_start_row(file_obj)
    
    try:
        df = pd.read_csv(file_obj, delimiter=delim, skiprows=skiprows, encoding='utf-8')
    except UnicodeDecodeError:
        file_obj.seek(0)
        df = pd.read_csv(file_obj, delimiter=delim, skiprows=skiprows, encoding='latin1')
        
    try:
        df = df.iloc[:, usecols]
        df.columns = colnames
    except IndexError:
        raise ValueError("Column index out of bounds.")
    return df

def validate_request(file_keys, form_keys):
    files = {}
    for key in file_keys:
        f = request.files.get(key)
        if not f or f.filename == '':
            return None, None, (jsonify({'error': f'Missing file: {key}'}), 400)
        files[key] = f

    params = {}
    for key in form_keys:
        val = request.form.get(key)
        try:
            params[key] = float(val) if (key != 'lastModified') else val
        except ValueError:
            params[key] = val
    
    # Read Switches
    params['createPPT'] = request.form.get('createPPT') == 'true'
    params['saveProject'] = request.form.get('saveProject') == 'true'
            
    return files, params, None

# --- Graph Helpers ---
def split_data_by_temp(df, mode='min'):
    idx = df['Temperature'].idxmin() if mode == 'min' else df['Temperature'].idxmax()
    return df.iloc[:idx + 1], df.iloc[idx + 1:]

def setup_graph(template, x_title, y_title, legend_txt, label_txt):
    graph = op.new_graph(template=template)
    layer = graph[0]
    layer.axis('x').title = x_title
    layer.axis('y').title = y_title
    if legend_txt: layer.label('Legend').text = legend_txt
    if label_txt: layer.label('Text').text = label_txt
    op.wait('s', 0.1)
    return graph

def add_plot(layer, wks, x_col, y_col, color_idx, name=None):
    plot = layer.add_plot(wks, colx=x_col, coly=y_col)
    plot.color = op.ocolor(color_idx)
    if name: plot.name = name
    op.wait('s', 0.1)
    return plot

# --- EXPORT LOGIC ---
def finalize_origin(should_save):
    if should_save:
        log_status("Saving Origin Project...")
        project_file = os.path.join(os.getcwd(), 'Origin_Project.opju')
        if os.path.exists(project_file):
            try:
                os.rename(project_file, project_file)
            except OSError:
                return "Project file locked. Skipping save."
        op.save(project_file)
    op.detach()
    return None

def export_graphs_to_pptx(graphs_list, filename, should_export):
    if not should_export:
        return None
    try:
        path = os.path.join(os.getcwd(), filename)
        if os.path.exists(path):
            try:
                os.rename(path, path)
            except OSError:
                return f"File {filename} is open. Close it to save."
        
        prs = Presentation()
        blank = prs.slide_layouts[6] 
        for graph in graphs_list:
            temp = os.path.join(os.getcwd(), f"temp_{id(graph)}.png")
            graph.save_fig(temp, type='png', width=1200)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(temp, Inches(0.5), Inches(1), height=Inches(5.5))
            try: os.remove(temp)
            except OSError: pass
        prs.save(path)
    except Exception as e:
        return str(e)
    return None

# --- Routes (All updated with flags) ---
@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/dewar', methods=['POST'])
def upload_dewar_file():
    with OriginContext():
        files, params, error = validate_request(['cooling', 'warming'], ['pressure', 'lastModified'])
        if error: return error
        fmt_date = params.get('lastModified', '')
        df_cool = load_data(files['cooling'], 3, [0, 3, 4], ['Temperature', 'R1', 'R2'])
        df_warm = load_data(files['warming'], 3, [0, 3, 4], ['Temperature', 'R1', 'R2'])
        wks_c = op.new_sheet('w', lname=f'CoolingData {params["pressure"]} GPa')
        wks_c.from_df(df_cool)
        wks_w = op.new_sheet('w', lname=f'WarmingData {params["pressure"]} GPa')
        wks_w.from_df(df_warm)
        graphs = [] 
        for i, ch in enumerate(['1', '2']):
            y_idx = 1 if ch == '1' else 2
            txt = f'{fmt_date}\nHg1223\nCh. {ch}\nPressure: {params["pressure"]} GPa'
            graph = setup_graph('Scatter', 'T (K)', 'R (Ω)', '\l(1) Cooling\n\l(2) Warming', txt)
            add_plot(graph[0], wks_c, 0, y_idx, 'blue')
            add_plot(graph[0], wks_w, 0, y_idx, 'red')
            graph[0].rescale()
            graphs.append(graph)
        
        err1 = export_graphs_to_pptx(graphs, f"Dewar_{params['pressure']}GPa.pptx", params['createPPT'])
        err2 = finalize_origin(params['saveProject'])
        msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
        return jsonify({'message': msg}), 200

def handle_single_file_resistance(req, skip, cols, split_mode):
    files, params, error = validate_request(['datafile'], ['pressure', 'lastModified'])
    if error: return error
    fmt_date = params.get('lastModified', '')
    df = load_data(files['datafile'], skip, cols, ['Temperature', 'R1', 'R2'])
    cool, warm = split_data_by_temp(df, split_mode)
    wks_c = op.new_sheet('w', lname=f'CoolingData {params["pressure"]} GPa')
    wks_c.from_df(cool)
    wks_w = op.new_sheet('w', lname=f'WarmingData {params["pressure"]} GPa')
    wks_w.from_df(warm)
    graphs = []
    for ch in ['1', '2']:
        y_idx = 1 if ch == '1' else 2
        txt = f'{fmt_date}\nHg1223\nCh. {ch}\nPressure: {params["pressure"]} GPa'
        graph = setup_graph('Scatter', 'T (K)', 'R (Ω)', '\l(1) Cooling\n\l(2) Warming', txt)
        add_plot(graph[0], wks_c, 0, y_idx, 'blue')
        add_plot(graph[0], wks_w, 0, y_idx, 'red')
        graph[0].rescale()
        graphs.append(graph)
    
    err1 = export_graphs_to_pptx(graphs, f"Resistance_{params['pressure']}GPa.pptx", params['createPPT'])
    err2 = finalize_origin(params['saveProject'])
    msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
    return jsonify({'message': msg}), 200

@app.route('/dewar_strip', methods=['POST'])
def upload_dewar_merged():
    with OriginContext(): return handle_single_file_resistance(request, 26, [0, 3, 4], 'min')

@app.route('/ppms', methods=['POST'])
def upload_ppms_file():
    with OriginContext(): return handle_single_file_resistance(request, -1, [3, 12, 13], 'min')

@app.route('/current_effect', methods=['POST'])
def current_effect():
    with OriginContext():
        files, params, error = validate_request(['datafile'], ['pressure', 'lastModified'])
        if error: return error
        fmt_date = params.get('lastModified', '')
        df = load_data(files['datafile'], 3, [0, 1, 2, 11], ['Temperature', 'R1', 'R2', 'Current'])
        df = df[pd.to_numeric(df['Temperature'], errors='coerce').notnull()]
        unique_currents = df['Current'].unique()
        wb = op.new_book('w', lname=f'CurrentData {params["pressure"]} GPa')
        graphs = []
        for ch_idx, ch_name in [(1, '1'), (2, '2')]:
            graph = setup_graph('Scatter', 'T (K)', 'R (Ω)', '', f'{fmt_date}\nHg1223\nCh. {ch_name}\n{params["pressure"]} GPa')
            legend = ''
            for i, curr in enumerate(unique_currents):
                sub = df[df['Current'] == curr]
                if sub.empty: continue
                wks = wb.add_sheet(f'Ch{ch_name}_{curr}mA')
                wks.from_df(sub[['Temperature', f'R{ch_name}']])
                add_plot(graph[0], wks, 0, 1, i)
                legend += f'\l({i+1}) {curr} A\n'
            graph[0].label('Legend').text = legend
            graph[0].rescale()
            graphs.append(graph)
        err1 = export_graphs_to_pptx(graphs, f"CurrentEffect_{params['pressure']}GPa.pptx", params['createPPT'])
        err2 = finalize_origin(params['saveProject'])
        msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
        return jsonify({'message': msg}), 200

@app.route('/ppms_magnetic', methods=['POST'])
def upload_ppms_magnetic_file():
    with OriginContext():
        files, params, error = validate_request(['datafile'], ['pressure', 'lastModified'])
        if error: return error
        df = load_data(files['datafile'], -1, [3, 4, 12, 13], ['Temperature', 'MagneticField', 'R1', 'R2'])
        unique_fields = df['MagneticField'].unique()
        wb = op.new_book('w', lname=f'MagneticFieldData {params["pressure"]} GPa')
        graphs = []
        fmt_date = params.get('lastModified', '')
        for ch_idx, ch_name in [(1, '1'), (2, '2')]:
            graph = setup_graph('Scatter', 'T (K)', 'R (Ω)', '', f'{fmt_date}\nCe\nCh. {ch_name}\n{params["pressure"]} GPa')
            legend = ''
            for i, field in enumerate(unique_fields):
                sub = df[df['MagneticField'] == field]
                if sub.empty: continue
                wks = wb.add_sheet(f'Field_{field}')
                wks.from_df(sub[['Temperature', f'R{ch_name}']])
                add_plot(graph[0], wks, 0, 1, i)
                legend += f'\l({i+1}) {round(field)/1000} T\n'
            graph[0].label('Legend').text = legend
            graph[0].rescale()
            graphs.append(graph)
        err1 = export_graphs_to_pptx(graphs, f"MagneticField_{params['pressure']}GPa.pptx", params['createPPT'])
        err2 = finalize_origin(params['saveProject'])
        msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
        return jsonify({'message': msg}), 200

@app.route('/ppms_heat_capacity', methods=['POST'])
def upload_ppms_heat_capacity_file():
    with OriginContext():
        files, params, error = validate_request(['datafile'], ['mass_heat_cap', 'lastModified'])
        if error: return error
        df = load_data(files['datafile'], -1, [7, 5, 9], ['Temperature', 'MagneticField', 'Heat capacity'])
        df['MagneticField'] = np.ceil(df['MagneticField'] / 10) * 10
        unique_fields = df['MagneticField'].unique()
        wb = op.new_book('w', lname=f'MagneticFieldData {params["mass_heat_cap"]} mg')
        fmt_date = params.get('lastModified', '')
        graph = setup_graph('Scatter', 'T (K)', 'Cp (mj/mole$\cdot$K)', '', f'{fmt_date}\n{params["mass_heat_cap"]} mg')
        legend = ''
        for i, field in enumerate(unique_fields):
            sub = df[df['MagneticField'] == field]
            if sub.empty: continue
            wks = wb.add_sheet(f'Field_{field}')
            wks.from_df(sub[['Temperature', 'Heat capacity']])
            add_plot(graph[0], wks, 0, 1, i)
            legend += f'\l({i+1}) {round(field)/1000} T\n'
        graph[0].label('Legend').text = legend
        graph[0].rescale()
        err1 = export_graphs_to_pptx([graph], f"HeatCapacity_{params['mass_heat_cap']}mg.pptx", params['createPPT'])
        err2 = finalize_origin(params['saveProject'])
        msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
        return jsonify({'message': msg}), 200

@app.route('/ppms_heat_capacity_cw', methods=['POST'])
def upload_ppms_heat_capacity_cw_file():
    with OriginContext():
        return handle_field_warming_cooling(request, skip=-1, cols=[7, 5, 9], col_names=['Temperature', 'Magnetic field', 'Heat capacity'], y_col='Heat capacity', y_label='Heat capacity (mj/mole$\cdot$K)', mass_key='mass', round_field=True)

@app.route('/mpms_magnetic', methods=['POST'])
def upload_mpms_magnetic_file():
    with OriginContext():
        return handle_field_warming_cooling(request, skip=-1, cols=[2, 3, 60], col_names=['Temperature', 'Magnetic field', 'Magnetic moment'], y_col='Magnetic moment', y_label='Magnetic moment (Oe)', mass_key='mass', round_field=False)

def handle_field_warming_cooling(req, skip, cols, col_names, y_col, y_label, mass_key, round_field):
    files, params, error = validate_request(['datafile'], [mass_key, 'lastModified'])
    if error: return error
    fmt_date = params.get('lastModified', '')
    df = load_data(files['datafile'], skip, cols, col_names)
    if round_field: df['Magnetic field'] = np.ceil(df['Magnetic field'] / 10) * 10
    field_col = 'Magnetic field'
    unique_fields = df[field_col].unique()
    wb = op.new_book('w', lname=f'Data_{params[mass_key]}mg')
    graph_warm = setup_graph('Scatter', 'T (K)', y_label, '', f'{fmt_date}\nZFC\nCe\nMass = {params[mass_key]}mg')
    graph_cool = setup_graph('Scatter', 'T (K)', y_label, '', f'{fmt_date}\nFC\nCe\nMass = {params[mass_key]}mg')
    leg_w, leg_c = '', ''
    for i, field in enumerate(unique_fields):
        sub = df[df[field_col] == field]
        if sub.empty: continue
        idx = sub['Temperature'].idxmax()
        warm_df, cool_df = sub.loc[:idx], sub.loc[idx+1:]
        wks_w = wb.add_sheet(f'Warming_{field}')
        wks_w.from_df(warm_df[['Temperature', y_col]])
        add_plot(graph_warm[0], wks_w, 0, 1, i, f'Warming_{field}')
        leg_w += f'\l({i+1}){field} Oe\n'
        wks_c = wb.add_sheet(f'Cooling_{field}')
        wks_c.from_df(cool_df[['Temperature', y_col]])
        add_plot(graph_cool[0], wks_c, 0, 1, i, f'Cooling_{field}')
        leg_c += f'\l({i+1}){field} Oe\n'
    graph_warm[0].label('Legend').text = leg_w
    graph_warm[0].rescale()
    graph_cool[0].label('Legend').text = leg_c
    graph_cool[0].rescale()
    err1 = export_graphs_to_pptx([graph_warm, graph_cool], f"CW_FieldData_{params[mass_key]}mg.pptx", params['createPPT'])
    err2 = finalize_origin(params['saveProject'])
    msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
    return jsonify({'message': msg}), 200

@app.route('/mpms', methods=['POST'])
def upload_mpms_file():
    with OriginContext():
        files, params, error = validate_request(['datafile'], ['magnetic_moment', 'lastModified'])
        if error: return error
        fmt_date = params.get('lastModified', '')
        df = load_data(files['datafile'], -1, [2, 60], ['Temperature', 'Magnetic moment'])
        warm, cool = split_data_by_temp(df, 'max')
        wks_w = op.new_sheet('w', lname=f'WarmingData {params["magnetic_moment"]} Oe')
        wks_w.from_df(warm)
        wks_c = op.new_sheet('w', lname=f'CoolingData {params["magnetic_moment"]} Oe')
        wks_c.from_df(cool)
        txt = f'{fmt_date}\nCe\nMagnetic field: {params["magnetic_moment"]} Oe'
        graph = setup_graph('Scatter', 'T (K)', 'Magnetic moment (Oe)', '\l(1) ZFC\n\l(2) FC', txt)
        add_plot(graph[0], wks_w, 0, 1, 'red')
        add_plot(graph[0], wks_c, 0, 1, 'blue')
        graph[0].rescale()
        err1 = export_graphs_to_pptx([graph], f"MPMS_{params['magnetic_moment']}Oe.pptx", params['createPPT'])
        err2 = finalize_origin(params['saveProject'])
        msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
        return jsonify({'message': msg}), 200

@app.route('/mpms_ac', methods=['POST'])
def upload_mpms_ac_file():
    with OriginContext():
        files, params, error = validate_request(['datafile'], ['mass_ac', 'MF_dc', 'MF_ac', 'lastModified'])
        if error: return error
        df = load_data(files['datafile'], -1, [2, 26, 21, 23], ['Temperature', 'Frequency', 'X_real', 'X_imag'])
        df = df.dropna()
        freqs = df['Frequency'].unique()
        if len(freqs) == 0: return jsonify({'error': 'No valid frequency data found'}), 400
        wb = op.new_book('w', lname=f"AC_Susceptibility_{params['mass_ac']}mg")
        graph_real = setup_graph('Scatter', 'Temperature (K)', "X' (emu/Oe)", '', '')
        graph_imag = setup_graph('Scatter', 'Temperature (K)', "X'' (emu/Oe)", '', '')
        legend = ''
        for i, freq in enumerate(freqs):
            sub = df[df['Frequency'] == freq]
            safe_freq = f"{freq:.2f}".replace('.', '_')
            wks = wb.add_sheet(f'Freq_{safe_freq}')
            wks.from_df(sub) 
            add_plot(graph_real[0], wks, 0, 2, i)
            add_plot(graph_imag[0], wks, 0, 3, i)
            legend += f'\l({i+1}) {freq:.1f} Hz\n'
        for graph in [graph_real, graph_imag]:
            graph[0].label('Legend').text = legend
            graph[0].rescale()
            op.wait('s', 0.1)
        err1 = export_graphs_to_pptx([graph_real, graph_imag], f"AC_Susceptibility_{params['mass_ac']}mg.pptx", params['createPPT'])
        err2 = finalize_origin(params['saveProject'])
        msg = "Processed successfully." if not (err1 or err2) else f"Done. Warnings: {err1 or ''} {err2 or ''}"
        return jsonify({'message': msg}), 200

# --- Force Close Splash Helper ---
def close_splash_forcefully():
    """
    Runs in a background thread. 
    Waits 3 seconds, then ruthlessly kills the splash screen.
    """
    time.sleep(3)
    if getattr(sys, 'frozen', False):
        try:
            import pyi_splash
            if pyi_splash.is_alive():
                pyi_splash.close()
        except:
            pass

def start_server():
    app.run(host='127.0.0.1', port=5000, threaded=False)

if __name__ == '__main__':
    # 1. Start the "Splash Killer" thread immediately
    threading.Thread(target=close_splash_forcefully, daemon=True).start()

    # 2. Start the Flask Server
    if os.environ.get('WERKZEUG_RUN_MAIN') != 'true':
        t = threading.Thread(target=start_server)
        t.daemon = True
        t.start()
        
        # 3. Start the GUI
        webview.create_window(
            "Lab Automation Dashboard", 
            "http://127.0.0.1:5000", 
            width=1000, 
            height=800
        )
        
        webview.start()